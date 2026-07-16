import math, pathlib, re
from datetime import datetime, timezone
import httpx
from bs4 import BeautifulSoup
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader
from sqlalchemy import delete, func, select
from sqlalchemy.orm import Session
from app.config import settings
from app.database import SessionLocal
from app.models import AgentFile, DocumentChunk, File, Folder

INDEXABLE_MIMES={"application/pdf","text/plain","text/markdown","text/csv","text/html","application/vnd.openxmlformats-officedocument.wordprocessingml.document","application/vnd.openxmlformats-officedocument.spreadsheetml.sheet","application/vnd.openxmlformats-officedocument.presentationml.presentation"}

def extract_sections(item:File)->list[tuple[str,str]]:
    path=pathlib.Path(item.path); mime=item.mime_type
    if mime=="application/pdf": return [(f"página {i+1}",page.extract_text() or "") for i,page in enumerate(PdfReader(path).pages)]
    if mime.endswith("wordprocessingml.document"):
        document=DocxDocument(path); return [("documento","\n".join(p.text for p in document.paragraphs))]
    if mime.endswith("spreadsheetml.sheet"):
        book=load_workbook(path,read_only=True,data_only=True); return [(f"aba {sheet.title}","\n".join(" | ".join("" if value is None else str(value) for value in row) for row in sheet.iter_rows(values_only=True))) for sheet in book.worksheets]
    if mime.endswith("presentationml.presentation"):
        deck=Presentation(path); return [(f"slide {i+1}","\n".join(shape.text for shape in slide.shapes if hasattr(shape,"text"))) for i,slide in enumerate(deck.slides)]
    raw=path.read_text(encoding="utf-8",errors="ignore")
    if mime=="text/html": raw=BeautifulSoup(raw,"html.parser").get_text("\n",strip=True)
    return [("arquivo",raw)]

def chunk_sections(sections:list[tuple[str,str]],target:int=1400,overlap:int=220)->list[tuple[str,str]]:
    chunks=[]
    for locator,text in sections:
        clean=re.sub(r"[ \t]+"," ",text); clean=re.sub(r"\n{3,}","\n\n",clean).strip()
        start=0
        while start<len(clean):
            end=min(start+target,len(clean))
            if end<len(clean):
                boundary=max(clean.rfind("\n",start+target//2,end),clean.rfind(". ",start+target//2,end))
                if boundary>start: end=boundary+1
            value=clean[start:end].strip()
            if len(value)>=40: chunks.append((locator,value))
            if end>=len(clean): break
            start=max(start+1,end-overlap)
    return chunks

def embed_texts(texts:list[str])->list[list[float]]:
    vectors=[]
    with httpx.Client(timeout=120) as client:
        for start in range(0,len(texts),32):
            response=client.post(f"{settings.deepinfra_base_url}/embeddings",json={"model":settings.embedding_ai_model,"input":texts[start:start+32]},headers={"Authorization":f"Bearer {settings.deepinfra_api_key}"})
            response.raise_for_status(); data=sorted(response.json()["data"],key=lambda item:item["index"])
            batch=[item["embedding"] for item in data]
            if any(len(vector)!=settings.embedding_dimensions for vector in batch): raise ValueError(f"O modelo {settings.embedding_ai_model} não retornou {settings.embedding_dimensions} dimensões")
            vectors.extend(batch)
    return vectors

def index_file(file_id:str)->dict:
    db=SessionLocal(); item=db.get(File,file_id)
    if not item: db.close(); return {"file_id":file_id,"status":"missing"}
    try:
        if item.mime_type not in INDEXABLE_MIMES: item.index_status="unsupported"; db.commit(); return {"file_id":file_id,"status":"unsupported"}
        if not settings.deepinfra_api_key: raise RuntimeError("DEEPINFRA_API_KEY não configurada")
        item.index_status="processing"; item.index_error=None; db.commit()
        sections=extract_sections(item); chunks=chunk_sections(sections)
        if not chunks: raise ValueError("Nenhum texto utilizável encontrado")
        vectors=embed_texts([content for _,content in chunks])
        db.execute(delete(DocumentChunk).where(DocumentChunk.file_id==item.id))
        for index,((locator,content),embedding) in enumerate(zip(chunks,vectors)):
            db.add(DocumentChunk(company_id=item.company_id,file_id=item.id,chunk_index=index,content=content,locator=locator,token_estimate=max(1,len(content)//4),embedding=embedding))
        item.extracted_text="\n\n".join(content for _,content in chunks)[:200000]; item.index_status="ready"; item.indexed_at=datetime.now(timezone.utc); db.commit()
        return {"file_id":file_id,"status":"ready","chunks":len(chunks)}
    except Exception as exc:
        db.rollback(); item=db.get(File,file_id)
        if item: item.index_status="failed"; item.index_error=str(exc)[:2000]; db.commit()
        raise
    finally: db.close()

def cosine_similarity(left:list[float],right:list[float])->float:
    numerator=sum(a*b for a,b in zip(left,right)); denominator=math.sqrt(sum(a*a for a in left))*math.sqrt(sum(b*b for b in right))
    return numerator/denominator if denominator else 0.0
def keyword_similarity(query:str,content:str)->float:
    terms={term for term in re.findall(r"[\w-]+",query.lower()) if len(term)>=3}
    if not terms: return 0.0
    haystack=content.lower(); return sum(term in haystack for term in terms)/len(terms)

def _allowed(file:File,folder:Folder|None,linked_agent_id:str|None,user_id:str,user_role:str,agent_id:str|None)->bool:
    if user_role in {"owner","admin","superadmin"} or file.user_id==user_id or agent_id and linked_agent_id==agent_id: return True
    if not folder or not folder.shared: return False
    permissions=folder.permissions or {}; denied=set(permissions.get("denied_user_ids",[])); users=set(permissions.get("user_ids",[])); roles=set(permissions.get("roles",[]))
    return user_id not in denied and (not users and not roles or user_id in users or user_role in roles)

def retrieve_chunks(db:Session,company_id:str,user_id:str,user_role:str,query_text:str,query_embedding:list[float],agent_id:str|None=None,limit:int=8):
    statement=select(DocumentChunk,File,Folder,AgentFile.agent_id).join(File,File.id==DocumentChunk.file_id).outerjoin(Folder,Folder.id==File.folder_id).outerjoin(AgentFile,AgentFile.file_id==File.id).where(DocumentChunk.company_id==company_id,File.index_status=="ready")
    if db.bind.dialect.name=="postgresql":
        distance=DocumentChunk.embedding.cosine_distance(query_embedding); keyword=func.ts_rank_cd(func.to_tsvector(DocumentChunk.content),func.plainto_tsquery(query_text)); rows=db.execute(statement.add_columns(distance,keyword).order_by(distance).limit(max(100,limit*10))).all()
        candidates=[]
        for chunk,file,folder,linked_agent,distance,keyword_rank in rows:
            semantic=max(0.0,1-float(distance)); lexical=min(1.0,float(keyword_rank or 0)*4); score=semantic*.8+lexical*.2
            if _allowed(file,folder,linked_agent,user_id,user_role,agent_id) and (semantic>=settings.rag_min_similarity or lexical>.15): candidates.append((chunk,file,score))
        candidates.sort(key=lambda row:row[2],reverse=True)
        seen=set(); return [row for row in candidates if not (row[0].id in seen or seen.add(row[0].id))][:limit]
    rows=db.execute(statement.limit(500)).all(); scored=[]
    for chunk,file,folder,linked_agent in rows:
        if not _allowed(file,folder,linked_agent,user_id,user_role,agent_id): continue
        semantic=cosine_similarity(chunk.embedding,query_embedding); lexical=keyword_similarity(query_text,chunk.content); score=semantic*.8+lexical*.2
        if semantic>=settings.rag_min_similarity or lexical>.15: scored.append((chunk,file,score))
    return sorted(scored,key=lambda row:row[2],reverse=True)[:limit]
